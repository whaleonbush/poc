# -*- coding: utf-8 -*-
"""외부 AI 4종 — 코딩·추론 역량 벤치마크 (중학생 눈높이, 그래프 중심)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

NAVY   = RGBColor(0x14, 0x2A, 0x4A)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
SKY    = RGBColor(0xE8, 0xF1, 0xFA)
GRAY   = RGBColor(0x44, 0x4A, 0x55)
LGRAY  = RGBColor(0x8A, 0x92, 0x9E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF2, 0x7E, 0x2E)
PEACH  = RGBColor(0xFE, 0xF1, 0xE6)
LINE   = RGBColor(0xD5, 0xDD, 0xE6)
LBLUE  = RGBColor(0x9F, 0xC4, 0xE8)
FONT   = "Apple SD Gothic Neo"

C_CLAUDE = RGBColor(0xD4, 0x77, 0x4A)
C_GPT    = RGBColor(0x10, 0xA3, 0x7A)
C_GEMINI = RGBColor(0x42, 0x85, 0xF4)
C_GROK   = RGBColor(0x6B, 0x7B, 0x8C)
MODEL_COLORS = [C_CLAUDE, C_GPT, C_GEMINI, C_GROK]
MODEL_SHORT = ["Claude\nOpus 4.8", "GPT-5.5", "Gemini\n3.1 Pro", "Grok 4.3"]
MODEL_PLAIN = ["Claude Opus 4.8", "GPT-5.5", "Gemini 3.1 Pro", "Grok 4.3"]

# 출처: Anthropic System Card · OpenAI/Gemini/xAI 공식 · Vellum/Artificial Analysis 교차
# * = 공식 미공개 → 독립 집계 근사치
DATA = {
    "swe_verified":  (88.6, 75.0, 80.6, 73.5),
    "swe_pro":       (69.2, 58.6, 54.2, None),
    "terminal":      (74.6, 78.2, 70.3, None),
    "osworld":       (83.4, 78.7, 75.0, None),
    "gpqa":          (93.6, 92.0, 94.3, 88.0),
    "hle_plain":     (49.8, 41.4, 44.4, 50.7),
    "hle_tools":     (57.9, 52.2, 51.4, None),
    "math500":       (96.1, 97.0, 95.6, 94.0),
    "aime":          (95.0, 97.0, 96.0, 96.0),
    "frontier_math": (49.0, 53.0, 51.0, 54.0),
}

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]
KICKER = "NC HW 개발팀  |  코딩·추론 역량 비교 (2026.6)"
TOTAL = 8

def add_slide(): return prs.slides.add_slide(BLANK)
def set_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, l, t, w, h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb

def header(slide, kicker, title):
    box(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    box(slide, 0, Inches(1.0), SW, Pt(3), fill=ACCENT)
    text(slide, Inches(0.55), Inches(0.13), Inches(12), Inches(0.3),
         [[(kicker, 11, LBLUE, True)]])
    text(slide, Inches(0.55), Inches(0.4), Inches(12.2), Inches(0.55),
         [[(title, 22, WHITE, True)]])

def pnum(slide, n):
    text(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
         [[(f"({n}/{TOTAL})", 10, LGRAY, False)]], align=PP_ALIGN.RIGHT)

def _style_chart(chart, title, y_max=None, show_legend=False, horizontal=False, title_size=12):
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    p = chart.chart_title.text_frame.paragraphs[0]
    p.font.size = Pt(title_size); p.font.bold = True; p.font.name = FONT; p.font.color.rgb = NAVY
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(8.5); chart.legend.font.name = FONT
    val_axis = chart.category_axis if horizontal else chart.value_axis
    cat_axis = chart.value_axis if horizontal else chart.category_axis
    val_axis.has_major_gridlines = True
    val_axis.tick_labels.font.size = Pt(9); val_axis.tick_labels.font.name = FONT
    cat_axis.tick_labels.font.size = Pt(9); cat_axis.tick_labels.font.name = FONT
    if y_max is not None:
        val_axis.maximum_scale = y_max; val_axis.minimum_scale = 0

def _color_series(series, color):
    fill = series.format.fill
    fill.solid(); fill.fore_color.rgb = color
    series.format.line.fill.background()

def _vals(key, include_grok=True):
    raw = DATA[key]
    if include_grok and raw[3] is not None:
        return raw
    return (raw[0], raw[1], raw[2], raw[3] or 0)

def _cats(include_grok=True):
    return MODEL_SHORT if include_grok else MODEL_SHORT[:3]

def _colors(include_grok=True):
    return MODEL_COLORS if include_grok else MODEL_COLORS[:3]

def bar_chart(slide, l, t, w, h, key, title, y_max=100, include_grok=True, sub=""):
    vals = _vals(key, include_grok)
    cats = _cats(include_grok)
    cols = _colors(include_grok)
    data = CategoryChartData()
    data.categories = cats
    data.add_series("", vals[:len(cats)])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, l, t, w, h, data
    ).chart
    _style_chart(chart, title + sub, y_max=y_max, horizontal=True, title_size=11.5)
    plot = chart.plots[0]
    series = plot.series[0]
    for idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid(); fill.fore_color.rgb = cols[idx]
        point.format.line.fill.background()
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9); plot.data_labels.font.name = FONT
    plot.data_labels.number_format = '0"%"'
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.category_axis.reverse_order = True
    return chart

def grouped_bar(slide, l, t, w, h, categories, series_list, title, y_max=100):
    data = CategoryChartData()
    data.categories = categories
    for name, vals in series_list:
        data.add_series(name, vals)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, l, t, w, h, data
    ).chart
    _style_chart(chart, title, y_max=y_max, show_legend=True, title_size=12)
    palette = [BLUE, C_CLAUDE, ACCENT, C_GPT, C_GEMINI]
    for i, s in enumerate(chart.plots[0].series):
        _color_series(s, palette[i % len(palette)])
    chart.value_axis.maximum_scale = y_max; chart.value_axis.minimum_scale = 0
    return chart

def legend_row(slide, l, t):
    items = [("Claude", C_CLAUDE), ("GPT", C_GPT), ("Gemini", C_GEMINI), ("Grok", C_GROK)]
    x = l
    for label, col in items:
        box(slide, x, t, Inches(0.16), Inches(0.16), fill=col)
        text(slide, x + Inches(0.22), t - Pt(1), Inches(1.2), Inches(0.2),
             [[(label, 9.5, GRAY, False)]])
        x += Inches(1.45)

def simple_box(slide, l, t, w, h, title, body, fill=SKY, accent=BLUE):
    box(slide, l, t, w, h, fill=fill, line=accent, line_w=1.5, round_=True)
    box(slide, l + Inches(0.12), t + Inches(0.12), Inches(0.45), Pt(3), fill=accent)
    text(slide, l + Inches(0.12), t + Inches(0.22), w - Inches(0.24), h - Inches(0.3),
         [[(title, 11.5, NAVY, True)], [(body, 10.5, GRAY, False)]], line_spacing=1.18)

def takeaway(slide, l, t, w, h, title, lines):
    box(slide, l, t, w, h, fill=NAVY, round_=True)
    runs = [[(title, 12, LBLUE, True)]]
    runs += [[(ln, 11, WHITE, False)] for ln in lines]
    text(slide, l + Inches(0.2), t, w - Inches(0.4), h, runs,
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=2)

# ── 1. 표지 ──
s = add_slide(); set_bg(s, NAVY)
text(s, Inches(0.95), Inches(2.0), Inches(11.6), Inches(1.3),
     [[("AI 코딩·추론 실력 비교", 38, WHITE, True)]])
text(s, Inches(0.97), Inches(3.15), Inches(11.6), Inches(0.55),
     [[("Claude · ChatGPT · Gemini · Grok — 시험 점수로 보는 그래프", 18, LBLUE, True)]])
box(s, Inches(0.97), Inches(3.85), Inches(4.6), Pt(4), fill=ACCENT)
text(s, Inches(0.97), Inches(4.15), Inches(11.6), Inches(1.0),
     [[("NC HW 개발팀  |  2026년 6월", 15, RGBColor(0xC9,0xD8,0xEA), False)],
      [("막대가 길수록 = 시험을 더 많이 맞힘 (100%가 만점)", 13, LGRAY, False)]],
     line_spacing=1.25)
pnum(s, 1)

# ── 2. 시험이 뭔지 설명 ──
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER, "먼저 알아두기 — '벤치마크'가 뭐예요?")
text(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.38),
     [[("벤치마크 = AI에게 내리는 표준 시험. 점수가 높을수록 그 종류 문제를 잘 푼다는 뜻입니다.", 12, GRAY, False)]])

cards = [
    ("코딩 시험", "프로그램 버그 고치기, 터미널 명령, 컴퓨터 조작", BLUE),
    ("추론 시험", "어려운 생각, 과학·수학 문제 풀기", ACCENT),
]
for i, (t, b, c) in enumerate(cards):
    simple_box(s, Inches(0.5 + i * 6.35), Inches(1.58), Inches(6.0), Inches(0.95), t, b, accent=c)

legend_row(s, Inches(0.5), Inches(2.72))
text(s, Inches(0.5), Inches(3.05), Inches(12.3), Inches(0.32),
     [[("아래 그래프는 같은 시험을 4개 AI가 친 점수입니다. (* 표시는 공식 미공개 → 근사치)", 10.5, LGRAY, False)]])

grouped_bar(s, Inches(0.5), Inches(3.42), Inches(12.3), Inches(2.55),
            ["코딩\n(평균)", "추론\n(평균)"],
            [
                ("Claude", (82.5, 61.3)),
                ("GPT", (77.5, 55.5)),
                ("Gemini", (75.0, 58.8)),
                ("Grok", (73.5, 52.4)),
            ],
            "영역별 평균 점수 — 어디가 강한지 큰 그림 (%)", y_max=100)

takeaway(s, Inches(0.5), Inches(6.18), Inches(12.35), Inches(0.88),
         "쉽게 말하면",
         ["코딩 = 코드·명령어 잘 다루는지  |  추론 = 머리 쓰는 문제·수학·과학 잘 푸는지",
          "한 AI가 모든 시험 1등은 아님 → 우리 팀 업무에 맞는 시험 점수를 보면 됨"])
pnum(s, 2)

# ── 3. 코딩 ① 프로그램 고치기 ──
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER, "코딩 ① — 실제 프로그램 고치기")
text(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.55),
     [[("SWE-bench = GitHub에 올라온 진짜 버그를 AI가 고칠 수 있는지 보는 시험", 11.5, GRAY, False)],
      [("Verified = 보통 난이도  |  Pro = 더 어려운 버전 (실무에 더 가까움)", 11.5, GRAY, False)]])

bar_chart(s, Inches(0.45), Inches(1.82), Inches(5.95), Inches(3.55),
          "swe_verified", "SWE-bench Verified (보통 난이도)", sub="\n버그 고치기 성공률")

bar_chart(s, Inches(6.85), Inches(1.82), Inches(5.95), Inches(3.55),
          "swe_pro", "SWE-bench Pro (어려운 버전)", include_grok=False,
          sub="\nGrok 미공개 — 3종만 비교")

takeaway(s, Inches(0.5), Inches(5.55), Inches(12.35), Inches(1.05),
         "코딩 ① 결론",
         ["보통·어려운 버그 고치기 모두 Claude 1위 (88.6% → 69.2%)",
          "GPT·Gemini도 50% 넘음 = 쓸 만한 수준  |  실무 코드 수정은 Claude가 가장 유리"])
pnum(s, 3)

# ── 4. 코딩 ② 터미널·컴퓨터 조작 ──
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER, "코딩 ② — 터미널·컴퓨터 직접 조작")
text(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.55),
     [[("Terminal-Bench = 명령어(터미널)로 서버·스크립트 조작 시험", 11.5, GRAY, False)],
      [("OSWorld = 마우스·키보드로 PC 화면을 직접 조작하는 시험 (자동화·RPA)", 11.5, GRAY, False)]])

bar_chart(s, Inches(0.45), Inches(1.82), Inches(5.95), Inches(3.55),
          "terminal", "Terminal-Bench 2.1 (터미널·명령어)", include_grok=False,
          sub="\n긴 명령 작업 성공률")

bar_chart(s, Inches(6.85), Inches(1.82), Inches(5.95), Inches(3.55),
          "osworld", "OSWorld-Verified (PC 직접 조작)", include_grok=False,
          sub="\nGemini·Grok * 근사치")

takeaway(s, Inches(0.5), Inches(5.55), Inches(12.35), Inches(1.05),
         "코딩 ② 결론",
         ["명령어·스크립트 작업 → GPT 1위 (78.2%)",
          "PC 화면 직접 조작 → Claude 1위 (83.4%)  |  HW팀: 터미널·자동화는 GPT, GUI 작업은 Claude"])
pnum(s, 4)

# ── 5. 추론 ① 과학·어려운 생각 ──
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER, "추론 ① — 과학 퀴즈 · 최난도 생각력")
text(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.55),
     [[("GPQA = 대학원 수준 과학(물리·화학·생물) 퀴즈 — 상위권은 90%대로 비슷함", 11.5, GRAY, False)],
      [("HLE = '인류 최후의 시험' — 도구 없이 혼자 푸는 초난도 종합 사고 시험", 11.5, GRAY, False)]])

bar_chart(s, Inches(0.45), Inches(1.82), Inches(5.95), Inches(3.55),
          "gpqa", "GPQA Diamond (대학원 과학 퀴즈)", sub="\n정답률")

bar_chart(s, Inches(6.85), Inches(1.82), Inches(5.95), Inches(3.55),
          "hle_plain", "HLE — 도구 없이 혼자 풀기", sub="\n* Grok 근사치")

takeaway(s, Inches(0.5), Inches(5.55), Inches(12.35), Inches(1.05),
         "추론 ① 결론",
         ["과학 퀴즈(GPQA)는 4종 모두 88~94% — 격차 작음 (거의 비슷)",
          "진짜 어려운 생각(HLE)에서는 Claude·Grok 상위 — '막히는 문제'는 Claude가 유리"])
pnum(s, 5)

# ── 6. 추론 ② 수학 ──
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER, "추론 ② — 수학 문제 (난이도별)")
text(s, Inches(0.5), Inches(1.12), Inches(12.3), Inches(0.55),
     [[("MATH-500 = 중·고등~대학 초급 수학  |  AIME = 올림피아드급  |  FrontierMath = 연구자도 어려운 최고난도", 11.5, GRAY, False)],
      [("쉬운 수학은 GPT·Gemini·Grok도 94%↑  |  어려울수록 점수 차이가 커짐", 11.5, GRAY, False)]])

grouped_bar(s, Inches(0.5), Inches(1.78), Inches(12.3), Inches(3.65),
            MODEL_SHORT,
            [
                ("MATH-500 (기본~중급)", DATA["math500"]),
                ("AIME 2026 (올림피아드급)", DATA["aime"]),
                ("FrontierMath (최고난도)", DATA["frontier_math"]),
            ],
            "수학 시험 3종 비교 — 막대가 길수록 정답률 높음 (%)", y_max=100)

takeaway(s, Inches(0.5), Inches(5.55), Inches(12.35), Inches(1.05),
         "추론 ② 결론",
         ["일반·올림피아드 수학 → GPT·Grok·Gemini 앞섬 (95~97%)",
          "최고난도 수학(FrontierMath) → Grok 54% 1위  |  HW팀: 계산·수식 많은 업무는 GPT, 초난도는 Grok도 검토"])
pnum(s, 6)

# ── 7. HLE 도구 사용 + 종합 ──
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER, "추론 ③ — 도구 쓰면 더 잘 푸는지 · 종합 정리")

bar_chart(s, Inches(0.45), Inches(1.55), Inches(6.0), Inches(3.2),
          "hle_tools", "HLE — 검색·계산 도구 사용 가능", include_grok=False,
          sub="\n도구 쓰면 점수 상승")

# 종합 winner cards
winners = [
    ("코딩\n버그 수정", "Claude", "SWE Pro 69%", C_CLAUDE),
    ("터미널\n명령어", "GPT", "Terminal 78%", C_GPT),
    ("PC\n직접 조작", "Claude", "OSWorld 83%", C_CLAUDE),
    ("어려운\n생각", "Claude", "HLE 50~58%", C_CLAUDE),
    ("과학\n퀴즈", "Gemini", "GPQA 94%", C_GEMINI),
    ("최고난도\n수학", "Grok", "Frontier 54%", C_GROK),
]
for i, (area, who, score, col) in enumerate(winners):
    cx = Inches(6.85) + (i % 3) * Inches(2.05)
    cy = Inches(1.55) + (i // 3) * Inches(2.05)
    box(s, cx, cy, Inches(1.9), Inches(1.85), fill=SKY, line=col, line_w=2, round_=True)
    text(s, cx + Inches(0.12), cy + Inches(0.15), Inches(1.65), Inches(1.55),
         [[(area, 10.5, NAVY, True)], [(who, 14, col, True)], [(score, 9.5, GRAY, False)]],
         align=PP_ALIGN.CENTER, line_spacing=1.15)

takeaway(s, Inches(0.5), Inches(5.72), Inches(12.35), Inches(0.95),
         "한 줄 요약",
         ["코딩·복잡한 생각 → Claude  |  터미널·수학 → GPT  |  과학 → Gemini  |  초난도 수학·저비용 → Grok",
          "시험마다 1등이 다름 → '우리 팀이 자주 하는 일'과 같은 시험 점수를 보면 선택이 쉬움"])
pnum(s, 7)

# ── 8. 출처 ──
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER, "점수는 어디서 가져왔나요?")
text(s, Inches(0.65), Inches(1.55), Inches(12), Inches(4.5),
     [[("공식 자료 (각 AI 회사가 직접 공개한 시험 점수)", 14, NAVY, True)],
      [("Anthropic Claude Opus 4.8 System Card  ·  OpenAI GPT-5.5  ·  Google Gemini 3.1 Pro  ·  xAI Grok 4.3", 11.5, GRAY, False)],
      [("", 8, GRAY, False)],
      [("독립 검증 (회사 말만 믿지 않고 다른 기관이 모아둔 자료)", 14, NAVY, True)],
      [("Artificial Analysis  ·  Vellum  ·  Weights & Biases ml-news", 11.5, GRAY, False)],
      [("", 8, GRAY, False)],
      [("꼭 기억할 점", 14, ACCENT, True)],
      [("① 같은 시험이어도 측정 방법이 다르면 점수가 달라질 수 있음", 11.5, GRAY, False)],
      [("② * 표시는 공식 미공개라 근사치 — 보고 시 '약' 또는 '추정'으로 말하면 됨", 11.5, GRAY, False)],
      [("③ 2~3% 차이는 오차 범위 — '비슷하다'고 봐도 됨", 11.5, GRAY, False)]],
     line_spacing=1.3)
pnum(s, 8)

prs.save("외부AI_4종_벤치마크_비교.pptx")
print("saved coding+reasoning charts:", TOTAL, "slides")
