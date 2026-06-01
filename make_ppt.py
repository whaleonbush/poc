# -*- coding: utf-8 -*-
"""HW 개발팀 AI 코딩 에이전트 활용 방안 보고 PPT 생성"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 색상/폰트 테마 ----------
NAVY   = RGBColor(0x14, 0x2A, 0x4A)   # 진한 네이비 (제목/포인트)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)   # 메인 블루
SKY    = RGBColor(0xE8, 0xF1, 0xFA)   # 연한 박스 배경
GRAY   = RGBColor(0x44, 0x4A, 0x55)   # 본문 텍스트
LGRAY  = RGBColor(0x8A, 0x92, 0x9E)   # 보조 텍스트
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF2, 0x7E, 0x2E)   # 강조 오렌지
LINE   = RGBColor(0xD5, 0xDD, 0xE6)

FONT = "Apple SD Gothic Neo"  # macOS 기본 한글 폰트 (PPT 열리는 환경 폰트로 자동 대체)

prs = Presentation()
prs.slide_width  = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    from pptx.enum.shapes import MSO_SHAPE
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp_type, l, t, w, h)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph = list of (text, size, color, bold) tuples"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb


def header(slide, kicker, title):
    """상단 제목 영역 (내지 공통)"""
    box(slide, 0, 0, SW, Inches(1.15), fill=NAVY)
    box(slide, 0, Inches(1.15), SW, Pt(3), fill=ACCENT)
    text(slide, Inches(0.6), Inches(0.16), Inches(12), Inches(0.35),
         [[(kicker, 12, RGBColor(0x9F,0xC4,0xE8), True)]])
    text(slide, Inches(0.6), Inches(0.45), Inches(12), Inches(0.6),
         [[(title, 26, WHITE, True)]])


def page_num(slide, n):
    text(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.3),
         [[(str(n), 11, LGRAY, False)]], align=PP_ALIGN.RIGHT)


def table_simple(slide, l, t, w, headers, rows, col_w, row_h=Inches(0.5),
                 head_fill=BLUE, fsize=12, hsize=12.5):
    """간단한 테이블 (도형 기반)"""
    x = l
    # header
    for j, htxt in enumerate(headers):
        cw = col_w[j]
        box(slide, x, t, cw, row_h, fill=head_fill, line=WHITE, line_w=1)
        text(slide, x, t, cw, row_h, [[(htxt, hsize, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cw
    # body
    y = t + row_h
    for ri, row in enumerate(rows):
        x = l
        rh = row_h
        fill = WHITE if ri % 2 == 0 else SKY
        for j, cell in enumerate(row):
            cw = col_w[j]
            box(slide, x, y, cw, rh, fill=fill, line=LINE, line_w=0.75)
            al = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            text(slide, x + Pt(4), y, cw - Pt(8), rh, [[(cell, fsize, GRAY, False)]],
                 align=al, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
            x += cw
        y += rh
    return y


# ============================================================
# Slide 1 — 표지
# ============================================================
s = add_slide(); set_bg(s, NAVY)
box(s, 0, Inches(4.05), SW, Pt(3.5), fill=ACCENT)
box(s, Inches(0.9), Inches(2.0), Inches(2.0), Pt(5), fill=BLUE)
text(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.6),
     [[("HW 개발팀 ", 40, WHITE, True), ("AI 코딩 에이전트", 40, RGBColor(0x9F,0xC4,0xE8), True)],
      [("활용 방안 보고", 40, WHITE, True)]], line_spacing=1.05)
text(s, Inches(0.92), Inches(4.35), Inches(11.5), Inches(0.6),
     [[("VS Code + Cline SR  ·  Gauss 사내 API  ·  Ollama 로컬 LLM", 16, RGBColor(0xC9,0xD8,0xEA), False)]])
text(s, Inches(0.92), Inches(6.5), Inches(11.5), Inches(0.6),
     [[("HW 개발팀  |  상무님 보고용  |  2026.06", 13, LGRAY, False)]])

# ============================================================
# Slide 2 — 한 장 요약
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "EXECUTIVE SUMMARY", "한 장 요약")
cards = [
    ("무엇을", "VS Code 기반 AI 코딩 에이전트 Cline SR을 사내 LLM(Gauss)·로컬 LLM(Ollama)과 연동하여 HW 개발 업무에 도입", BLUE),
    ("왜", "코드·설계자산을 외부로 유출하지 않는 사내망·로컬 환경에서 AI 활용 → 보안 위험 없이 생산성 향상", ACCENT),
    ("핵심 효과", "펌웨어·테스트·문서·리뷰 등 반복·정형 업무 자동화 → 개발 리드타임 단축 및 검증 품질 향상", BLUE),
    ("제언", "저위험·고효용 업무부터 단계적 시범 적용(PoC) 후 팀 전체로 확산", ACCENT),
]
top = Inches(1.5); ch = Inches(1.25); gap = Inches(0.12)
for i, (k, v, c) in enumerate(cards):
    y = top + i * (ch + gap)
    box(s, Inches(0.6), y, Pt(7), ch, fill=c)
    box(s, Inches(0.7), y, Inches(12.0), ch, fill=SKY if i % 2 == 0 else WHITE,
        line=LINE, line_w=1)
    text(s, Inches(0.95), y, Inches(2.2), ch, [[(k, 18, c, True)]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(3.2), y, Inches(9.3), ch, [[(v, 14.5, GRAY, False)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
page_num(s, 2)

# ============================================================
# Slide 3 — 추진 배경
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "BACKGROUND", "추진 배경")
items = [
    ("AI 코딩 에이전트 확산", "SW 직군 중심으로 빠르게 확산 중이나 HW 개발팀은 활용도가 낮음"),
    ("HW팀도 코드성 업무 多", "펌웨어·검증 스크립트·문서 등 코드성 업무 비중이 상당함"),
    ("외부 AI는 사용 제한", "ChatGPT 등 외부 서비스는 소스코드·설계자산 유출 우려로 사내 사용 제한"),
    ("사내 도구는 이미 안전", "Gauss(사내 API)·Ollama(로컬)는 외부 노출 없이 즉시 활용 가능"),
]
top = Inches(1.7)
for i, (t1, t2) in enumerate(items):
    y = top + i * Inches(1.15)
    box(s, Inches(0.7), y, Inches(0.55), Inches(0.55), fill=BLUE, round_=True)
    text(s, Inches(0.7), y, Inches(0.55), Inches(0.55), [[(str(i+1), 18, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(1.5), y - Inches(0.05), Inches(11), Inches(0.45),
         [[(t1, 17, NAVY, True)]])
    text(s, Inches(1.5), y + Inches(0.42), Inches(11), Inches(0.5),
         [[(t2, 14, GRAY, False)]])
page_num(s, 3)

# ============================================================
# Slide 4 — 도구 구성 및 동작 원리
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "ARCHITECTURE", "도구 구성 및 동작 원리")
# 좌: 구성요소 / 우: 개념도
text(s, Inches(0.6), Inches(1.35), Inches(6), Inches(0.4),
     [[("구성 요소", 16, NAVY, True)]])
comp = [
    ("VS Code", "개발자 표준 IDE (코드 편집기)"),
    ("Cline SR", "VS Code 내 AI 코딩 에이전트 — 코드 읽기/작성/수정·명령 실행을 자동 수행"),
    ("Gauss", "에이전트의 두뇌, 사내 LLM API (사내망, 고성능)"),
    ("Ollama", "PC에서 직접 구동되는 로컬 LLM (완전 오프라인)"),
]
y = Inches(1.85)
for name, desc in comp:
    box(s, Inches(0.6), y, Inches(6.0), Inches(1.05), fill=SKY, line=LINE, line_w=1, round_=True)
    text(s, Inches(0.85), y + Inches(0.13), Inches(5.5), Inches(0.4), [[(name, 15, BLUE, True)]])
    text(s, Inches(0.85), y + Inches(0.5), Inches(5.5), Inches(0.5), [[(desc, 12.5, GRAY, False)]],
         line_spacing=0.95)
    y += Inches(1.2)

# 우측 개념도
rx = Inches(7.0)
text(s, rx, Inches(1.35), Inches(6), Inches(0.4), [[("동작 개념도", 16, NAVY, True)]])
box(s, rx, Inches(1.95), Inches(5.7), Inches(0.8), fill=NAVY, round_=True)
text(s, rx, Inches(1.95), Inches(5.7), Inches(0.8), [[("개발자", 15, WHITE, True),(" · VS Code + Cline SR", 13, RGBColor(0xC9,0xD8,0xEA), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 두 갈래
box(s, rx, Inches(3.2), Inches(2.75), Inches(1.1), fill=BLUE, round_=True)
text(s, rx, Inches(3.2), Inches(2.75), Inches(1.1),
     [[("Gauss 사내 API", 14, WHITE, True)],[("사내망·고성능·팀 공용", 11.5, RGBColor(0xD8,0xE6,0xF4), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
box(s, rx + Inches(2.95), Inches(3.2), Inches(2.75), Inches(1.1), fill=BLUE, round_=True)
text(s, rx + Inches(2.95), Inches(3.2), Inches(2.75), Inches(1.1),
     [[("Ollama 로컬 LLM", 14, WHITE, True)],[("PC 내 완전 오프라인", 11.5, RGBColor(0xD8,0xE6,0xF4), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
# 보안 강조 박스
box(s, rx, Inches(4.65), Inches(5.7), Inches(1.15), fill=RGBColor(0xFE,0xF1,0xE6), line=ACCENT, line_w=1.5, round_=True)
text(s, rx + Inches(0.2), Inches(4.65), Inches(5.3), Inches(1.15),
     [[("\U0001F512  핵심 보안 포인트", 14, ACCENT, True)],
      [("입력한 코드·결과물이 사내망/로컬 PC 밖으로", 13, GRAY, False)],
      [("전혀 나가지 않음 → 외부 유출 위험 없음", 13, GRAY, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=2)
page_num(s, 4)

# ============================================================
# Slide 5 — Gauss vs Ollama
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "ARCHITECTURE", "Gauss vs Ollama — 언제 무엇을 쓰나")
table_simple(
    s, Inches(0.7), Inches(1.7), Inches(11.9),
    ["구분", "Gauss (사내 API)", "Ollama (로컬)"],
    [
        ["성능", "높음 (대형 모델)", "PC 사양에 의존 (상대적 제한)"],
        ["보안 수준", "사내망 내 안전", "최고 (오프라인·외부 차단)"],
        ["적합 업무", "일반 코드 작성·리뷰·문서화 등 대부분", "망분리 환경, 최고 민감도 자료 처리"],
        ["비용/자원", "공용 서버 자원", "개인 PC 자원"],
    ],
    col_w=[Inches(2.3), Inches(4.8), Inches(4.8)], row_h=Inches(0.95), fsize=13.5
)
box(s, Inches(0.7), Inches(6.5), Inches(11.9), Inches(0.7), fill=SKY, line=LINE, line_w=1, round_=True)
text(s, Inches(0.95), Inches(6.5), Inches(11.4), Inches(0.7),
     [[("\U0001F4A1 운영 원칙: ", 13.5, BLUE, True),("일반 업무는 Gauss, 망분리·최고민감 업무만 Ollama로 분리 운영", 13.5, GRAY, False)]],
     anchor=MSO_ANCHOR.MIDDLE)
page_num(s, 5)

# ============================================================
# Slide 6 — 활용 방안 (단기 Quick-win)
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "USE CASES \u2460", "HW 개발팀 활용 방안 — 우선 적용 (단기)")
table_simple(
    s, Inches(0.55), Inches(1.65), Inches(12.2),
    ["활용 방안", "내용", "기대 효과"],
    [
        ["펌웨어/드라이버 작성 보조", "C/C++ 드라이버·HAL·통신(I2C·SPI·UART·CAN) 코드 초안 생성", "작성 시간 단축"],
        ["레지스터 맵 자동화", "데이터시트의 레지스터 정의를 헤더(.h) 코드로 변환", "수기 입력 오류 감소"],
        ["테스트 자동화 스크립트", "계측기 제어(SCPI)·측정 자동화 Python 스크립트 생성", "검증 반복 업무 자동화"],
        ["측정 데이터 분석·시각화", "로그·측정 데이터 파싱, 그래프/리포트 생성 코드", "분석 시간 단축"],
        ["문서화 자동화", "코드 주석·README·설계/테스트 리포트 초안 작성", "문서 작성 부담 경감"],
    ],
    col_w=[Inches(3.3), Inches(6.0), Inches(2.9)], row_h=Inches(0.92), fsize=12.5, hsize=13
)
page_num(s, 6)

# ============================================================
# Slide 7 — 활용 방안 (중기 확산)
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "USE CASES \u2461", "HW 개발팀 활용 방안 — 확산 적용 (중기)")
table_simple(
    s, Inches(0.55), Inches(1.65), Inches(12.2),
    ["활용 방안", "내용", "기대 효과"],
    [
        ["코드 리뷰 보조", "코딩 규칙(MISRA-C) 점검, 잠재 버그·예외처리 누락 지적", "리뷰 품질·속도 향상"],
        ["레거시 코드 이해", "오래된/인수인계 펌웨어 코드 분석·설명", "온보딩 가속"],
        ["디버깅 보조", "에러 로그·증상 기반 원인 후보 및 수정안 제시", "문제 해결 시간 단축"],
        ["빌드/CI 스크립트", "Makefile·CMake·빌드 자동화 스크립트 작성·개선", "빌드 환경 정비"],
        ["HDL 테스트벤치 보조", "Verilog/VHDL testbench 초안 생성 (검토 필수)", "검증 환경 구축 보조"],
    ],
    col_w=[Inches(3.3), Inches(6.0), Inches(2.9)], row_h=Inches(0.92), fsize=12.5, hsize=13
)
page_num(s, 7)

# ============================================================
# Slide 8 — 활용 시나리오 예시
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "SCENARIOS", "활용 시나리오 예시")
scn = [
    ("드라이버 개발", "\u201C데이터시트의 레지스터 표 기반으로 센서 초기화 드라이버와 헤더를 작성해줘\u201D",
     "Cline SR이 초안 생성 → 개발자가 검토·수정·실측 검증"),
    ("테스트 자동화", "\u201C부하를 단계별로 바꾸며 전압·전류를 기록하는 효율 측정 스크립트를 만들어줘\u201D",
     "측정 자동화 스크립트 즉시 생성 → 반복 측정에 재사용"),
    ("문서화", "\u201C이 펌웨어 모듈의 동작을 설명하는 설계 문서 초안을 작성해줘\u201D",
     "코드 분석 후 설계 문서 초안 자동 작성"),
]
top = Inches(1.6)
for i, (tag, q, a) in enumerate(scn):
    y = top + i * Inches(1.75)
    box(s, Inches(0.7), y, Inches(11.9), Inches(1.55), fill=WHITE, line=LINE, line_w=1.2, round_=True)
    box(s, Inches(0.7), y, Inches(2.4), Inches(1.55), fill=NAVY, round_=True)
    text(s, Inches(0.7), y, Inches(2.4), Inches(1.55), [[("예시 "+str(i+1), 12, RGBColor(0x9F,0xC4,0xE8), True)],[(tag, 16, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    text(s, Inches(3.35), y + Inches(0.2), Inches(9.0), Inches(0.55),
         [[(q, 14, BLUE, True)]], line_spacing=1.0)
    text(s, Inches(3.35), y + Inches(0.85), Inches(9.0), Inches(0.55),
         [[("\u2192 ", 13, ACCENT, True),(a, 13, GRAY, False)]], line_spacing=1.0)
page_num(s, 8)

# ============================================================
# Slide 9 — 도입 로드맵
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "ROADMAP", "도입 로드맵")
steps = [
    ("1단계 · 준비", "2~3주", "환경 세팅(VS Code+Cline SR), Gauss/Ollama 연동, 보안 가이드 숙지", BLUE),
    ("2단계 · PoC", "4~6주", "Quick-win 과제 1~2건 시범 적용, 효과 측정", ACCENT),
    ("3단계 · 확산", "이후", "팀 표준 워크플로우 정립, 활용 사례·프롬프트 가이드 공유", BLUE),
]
cw = Inches(3.85); gap = Inches(0.25); left = Inches(0.7); top = Inches(2.0); ch = Inches(2.9)
for i, (title_, dur, desc, c) in enumerate(steps):
    x = left + i * (cw + gap)
    box(s, x, top, cw, ch, fill=WHITE, line=c, line_w=2, round_=True)
    box(s, x, top, cw, Inches(0.85), fill=c, round_=True)
    text(s, x, top, cw, Inches(0.85), [[(title_, 18, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, x, top + Inches(1.05), cw, Inches(0.5), [[(dur, 16, c, True)]],
         align=PP_ALIGN.CENTER)
    text(s, x + Inches(0.25), top + Inches(1.65), cw - Inches(0.5), Inches(1.1),
         [[(desc, 13, GRAY, False)]], align=PP_ALIGN.CENTER, line_spacing=1.1)
    if i < 2:
        text(s, x + cw - Inches(0.05), top + Inches(1.1), Inches(0.5), Inches(0.6),
             [[("\u25B6", 18, LGRAY, True)]], align=PP_ALIGN.CENTER)
box(s, Inches(0.7), Inches(5.5), Inches(11.85), Inches(1.0), fill=SKY, line=LINE, line_w=1, round_=True)
text(s, Inches(0.95), Inches(5.5), Inches(11.4), Inches(1.0),
     [[("\U0001F4CC PoC 운영 방안", 13.5, BLUE, True)],
      [("자원자 2~3명 + 명확한 과제로 시작 / 지표 예: 작성·문서화 소요시간, 리뷰 지적 건수, 스크립트 재사용률", 13, GRAY, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05, space_after=2)
page_num(s, 9)

# ============================================================
# Slide 10 — 기대 효과
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "EXPECTED BENEFITS", "기대 효과")
benefits = [
    ("생산성", "반복·정형 코드/문서 작성 시간 단축 → 설계·검증 등 고부가 업무에 집중"),
    ("품질", "코드 리뷰·규칙 점검 보조로 결함 조기 발견, 문서화율 향상"),
    ("지식 자산화", "레거시 코드 설명·문서 자동화로 속인화 위험 완화, 온보딩 가속"),
    ("보안 양립", "외부 유출 없이 AI 활용 → 생산성과 보안을 동시 확보"),
]
top = Inches(1.6); bw = Inches(5.85); bh = Inches(2.35); gapx = Inches(0.2); gapy = Inches(0.25)
for i, (k, v) in enumerate(benefits):
    r, cidx = divmod(i, 2)
    x = Inches(0.7) + cidx * (bw + gapx)
    y = top + r * (bh + gapy)
    box(s, x, y, bw, bh, fill=SKY, line=LINE, line_w=1, round_=True)
    box(s, x, y, bw, Pt(6), fill=ACCENT if i % 2 else BLUE)
    text(s, x + Inches(0.35), y + Inches(0.35), bw - Inches(0.7), Inches(0.6),
         [[(k, 22, NAVY, True)]])
    text(s, x + Inches(0.35), y + Inches(1.15), bw - Inches(0.7), Inches(1.0),
         [[(v, 15, GRAY, False)]], line_spacing=1.15)
page_num(s, 10)

# ============================================================
# Slide 11 — 리스크 및 고려사항
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "RISKS & MITIGATION", "리스크 및 고려사항")
table_simple(
    s, Inches(0.55), Inches(1.6), Inches(12.2),
    ["항목", "내용", "대응 방안"],
    [
        ["결과물 정확성", "AI 생성 코드는 오류·미흡 가능성 존재", "개발자 검토·실측 검증 후 적용 (AI는 보조)"],
        ["HW 특화 한계", "회로·타이밍 등 물리 영역은 직접 판단 불가", "코드·문서·분석 등 SW성 업무 중심 활용"],
        ["로컬 성능 제약", "Ollama는 PC 사양에 따라 성능 제한", "일반 업무는 Gauss, 고민감 업무만 Ollama"],
        ["활용 역량 편차", "사용법·프롬프트 숙련도 차이", "사내 가이드·사례 공유, 짧은 내부 교육"],
        ["보안 준수", "도구 사용 범위·반입 규정 준수 필요", "사내 보안 정책 범위 내 사용 원칙"],
    ],
    col_w=[Inches(2.6), Inches(4.9), Inches(4.7)], row_h=Inches(0.9), fsize=12.5, hsize=13
)
page_num(s, 11)

# ============================================================
# Slide 12 — 결론 및 제언
# ============================================================
s = add_slide(); set_bg(s, NAVY)
box(s, 0, Inches(2.7), SW, Pt(3.5), fill=ACCENT)
text(s, Inches(0.9), Inches(0.9), Inches(11.5), Inches(0.6),
     [[("결론 및 제언", 30, WHITE, True)]])
concl = [
    "VS Code + Cline SR + Gauss/Ollama는 보안 위험 없이 즉시 활용 가능한 사내 AI 코딩 환경",
    "HW 개발팀의 펌웨어·테스트·문서·리뷰 업무에 적용 가능, 특히 반복·정형 업무 자동화 효과 큼",
    "제언: 저위험·고효용 과제 중심의 단계적 PoC 후 팀 확산을 건의",
]
y = Inches(3.3)
for c in concl:
    box(s, Inches(0.9), y + Inches(0.07), Inches(0.16), Inches(0.16), fill=ACCENT)
    text(s, Inches(1.3), y - Inches(0.05), Inches(11), Inches(0.7),
         [[(c, 17, RGBColor(0xE8,0xF1,0xFA), False)]], line_spacing=1.1)
    y += Inches(0.95)
text(s, Inches(0.9), Inches(6.7), Inches(11.5), Inches(0.5),
     [[("※ 모든 AI 결과물은 개발자 검토·검증을 거쳐 적용함을 원칙으로 함", 12, LGRAY, False)]])

prs.save("HW개발팀_AI코딩에이전트_활용방안_보고.pptx")
print("saved:", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
