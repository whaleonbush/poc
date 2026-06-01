# -*- coding: utf-8 -*-
"""외부 AI 3종(ChatGPT·Gemini·Claude) 비교 조사 PPT"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x14, 0x2A, 0x4A)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
SKY    = RGBColor(0xE8, 0xF1, 0xFA)
GRAY   = RGBColor(0x44, 0x4A, 0x55)
LGRAY  = RGBColor(0x8A, 0x92, 0x9E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF2, 0x7E, 0x2E)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PEACH  = RGBColor(0xFE, 0xF1, 0xE6)
LINE   = RGBColor(0xD5, 0xDD, 0xE6)
FONT   = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide(): return prs.slides.add_slide(BLANK)
def set_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, l, t, w, h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb

def header(slide, kicker, title):
    box(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    box(slide, 0, Inches(1.0), SW, Pt(3), fill=ACCENT)
    text(slide, Inches(0.55), Inches(0.13), Inches(12), Inches(0.3),
         [[(kicker, 11, RGBColor(0x9F,0xC4,0xE8), True)]])
    text(slide, Inches(0.55), Inches(0.4), Inches(12.2), Inches(0.55),
         [[(title, 23, WHITE, True)]])

def section(slide, l, t, w, label, color=BLUE):
    box(slide, l, t, Pt(5), Inches(0.32), fill=color)
    text(slide, l + Inches(0.12), t - Inches(0.02), w, Inches(0.36),
         [[(label, 14.5, NAVY, True)]])

def pnum(slide, n):
    text(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
         [[(str(n), 10, LGRAY, False)]], align=PP_ALIGN.RIGHT)

def table_simple(slide, l, t, headers, rows, col_w, row_h, head_fill=BLUE, fsize=10.5, hsize=11):
    x = l
    for j, htxt in enumerate(headers):
        cw = col_w[j]
        box(slide, x, t, cw, row_h, fill=head_fill, line=WHITE, line_w=1)
        text(slide, x, t, cw, row_h, [[(htxt, hsize, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cw
    y = t + row_h
    for ri, row in enumerate(rows):
        x = l; fill = WHITE if ri % 2 == 0 else SKY
        for j, cell in enumerate(row):
            cw = col_w[j]
            box(slide, x, y, cw, row_h, fill=fill, line=LINE, line_w=0.75)
            al = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            text(slide, x + Pt(3), y, cw - Pt(6), row_h, [[(cell, fsize, GRAY, False)]],
                 align=al, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.92)
            x += cw
        y += row_h
    return y

def bullet_card(slide, l, t, w, h, title, bullets, title_color=BLUE):
    box(slide, l, t, w, h, fill=SKY, line=LINE, line_w=1, round_=True)
    box(slide, l, t, Pt(6), h, fill=title_color)
    text(slide, l + Inches(0.2), t + Inches(0.1), w - Inches(0.4), Inches(0.35),
         [[(title, 13, title_color, True)]])
    yy = t + Inches(0.48)
    for b in bullets:
        text(slide, l + Inches(0.25), yy, w - Inches(0.45), Inches(0.38),
             [[("\u2022 " + b, 11, GRAY, False)]], line_spacing=1.0)
        yy += Inches(0.4)

# Slide 1 — 표지
s = add_slide(); set_bg(s, NAVY)
box(s, 0, Inches(3.9), SW, Pt(4), fill=ACCENT)
text(s, Inches(0.9), Inches(1.8), Inches(11.5), Inches(1.5),
     [[("외부 AI 3종 비교 조사", 38, WHITE, True)],
      [("ChatGPT  \u00b7  Gemini  \u00b7  Claude", 22, RGBColor(0x9F,0xC4,0xE8), True)]],
     line_spacing=1.1)
text(s, Inches(0.92), Inches(4.2), Inches(11.5), Inches(0.8),
     [[("삼성전자 스마트폰 HW 개발팀  |  2026년 6월 외부 AI 도입 검토", 15, RGBColor(0xC9,0xD8,0xEA), False)],
      [("공식 문서 + 웹 사용기 교차 검증", 13, LGRAY, False)]])
pnum(s, 1)

# Slide 2 — 보고 목적
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (1/16)", "보고 목적 및 평가 기준")
section(s, Inches(0.55), Inches(1.2), Inches(6), "왜 지금 비교하나")
text(s, Inches(0.7), Inches(1.65), Inches(12), Inches(1.0),
     [[("2026년 6월부터 외부 AI 도입을 검토합니다. 도입 전 ChatGPT·Gemini·Claude 각각의", 13, GRAY, False)],
      [("강점과 차이를 알아두어, 우리 팀 업무에 맞는 도구를 고르기 위한 조사입니다.", 13, GRAY, False)]],
     line_spacing=1.12)
section(s, Inches(0.55), Inches(2.85), Inches(6), "평가 기준 (6가지)")
criteria = ["생산성 (코딩·문서·반복 업무)", "에이전트 (자동으로 여러 단계 수행)",
            "문서·지식 (긴 사양서·검색·OCR)", "멀티모달 (이미지·PDF 등)",
            "엔터프라이즈 (관리·보안·연동)", "비용·속도 (API·사용 제한)"]
y = Inches(3.3)
for c in criteria:
    box(s, Inches(0.7), y, Inches(5.8), Inches(0.48), fill=SKY, line=LINE, line_w=0.8, round_=True)
    text(s, Inches(0.9), y, Inches(5.5), Inches(0.48), [[("\u2022 " + c, 12, GRAY, False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)
box(s, Inches(7.0), Inches(2.85), Inches(5.75), Inches(3.5), fill=PEACH, line=ACCENT, line_w=1.2, round_=True)
text(s, Inches(7.25), Inches(3.0), Inches(5.3), Inches(3.0),
     [[("조사 방법", 14, ACCENT, True)],
      [("1차: 각 회사 공식 문서", 12, GRAY, False)],
      [("2차: 웹 사용기·비교 리뷰", 12, GRAY, False)],
      [("(참고로 표기)", 11, LGRAY, False)],
      [("", 8, GRAY, False)],
      [("주의: 보안·라이선스는", 12, GRAY, False)],
      [("별도 법무·보안 검토 필요", 12, RED, True)]],
     line_spacing=1.08, space_after=2)
pnum(s, 2)

# Slide 3 — 3종 한눈에
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (2/16)", "3종 한눈에 보기")
table_simple(
    s, Inches(0.55), Inches(1.55),
    ["AI", "한 줄 포지션", "강점 (공식·참고)", "주의"],
    [
        ["ChatGPT", "만능형 올라운더", "추론·코딩·AgentKit·MCP·도구 생태계 최대", "범용이라 업무별 최적은 다를 수 있음"],
        ["Gemini", "Google 연동·긴 문서", "3.5 Flash 에이전트·코딩·검색·멀티모달·속도", "Google 생태계 의존"],
        ["Claude", "코딩·장기 에이전트", "Opus/Sonnet·1M 맥락·Claude Code·구조화 출력", "창의·범용은 ChatGPT 대비 평가 분산"],
    ],
    col_w=[Inches(1.35), Inches(2.2), Inches(5.5), Inches(3.2)],
    row_h=Inches(0.95), fsize=11, hsize=11.5
)
box(s, Inches(0.55), Inches(5.0), Inches(12.25), Inches(1.5), fill=NAVY, round_=True)
text(s, Inches(0.8), Inches(5.15), Inches(11.8), Inches(1.2),
     [[("핵심 메시지", 14, RGBColor(0x9F,0xC4,0xE8), True)],
      [("세 AI 모두 '최고 하나'가 아니라, 업무 유형별로 맞는 도구가 다릅니다. 실무에서는 복수 도입·역할 분담도 흔합니다.", 12.5, WHITE, False)]],
     line_spacing=1.1, space_after=2)
pnum(s, 3)

# Slide 4 — ChatGPT
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (3/16)", "ChatGPT (OpenAI)")
bullet_card(s, Inches(0.55), Inches(1.35), Inches(6.0), Inches(2.5), "공식 강점",
    ["GPT-5.5: 복잡 추론·코딩용 플래그십 (맥락 약 100만 토큰)",
     "Responses API: 추론·도구·웹/파일 검색·Computer Use",
     "AgentKit: 에이전트 설계·ChatKit·MCP·Evals",
     "이미지·음성 등 전용 모델 라인업"], BLUE)
bullet_card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(2.5), "사용자 평가 (참고)",
    ["만능형: DevOps·인프라·프로토타입에 강점",
     "도구·플러그인 생태계가 가장 넓음",
     "엔터프라이즈 도입 사례·지원 체계 풍부"], ACCENT)
text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.4), [[("HW 예시", 13, NAVY, True)]])
text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(1.2),
     [[("\u2022 측정 자동화 스크립트 초안, CI/CD 로그 분석, 회의록·시험 보고서 초안", 12, GRAY, False)],
      [("\u2022 MCP로 사내 도구 연동 검토 시 후보 (보안 승인 후)", 12, GRAY, False)]],
     line_spacing=1.1)
pnum(s, 4)

# Slide 5 — Gemini
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (4/16)", "Gemini (Google)")
bullet_card(s, Inches(0.55), Inches(1.35), Inches(6.0), Inches(2.5), "공식 강점",
    ["Gemini 3.5 Flash: 에이전트·코딩·멀티모달 (2026.5 GA)",
     "Google Search·File Search·URL Context 그라운딩",
     "Managed Agents, Antigravity 개발 플랫폼",
     "100+ 페이지 문서 추론 사례 (금융 등)"], BLUE)
bullet_card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(2.5), "사용자 평가 (참고)",
    ["긴 코드베이스·사양서 한 번에 분석",
     "검색·최신 정보·비용 효율 강점",
     "보안·논리 분석 맥락 이해 평가"], ACCENT)
text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.4), [[("HW 예시", 13, NAVY, True)]])
text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(1.2),
     [[("\u2022 센서·PMIC 데이터시트(수백 페이지)에서 설정값 찾기", 12, GRAY, False)],
      [("\u2022 스캔 PDF·회로도 캡처 OCR 후 텍스트 정리 (공식 OCR 사례)", 12, GRAY, False)]],
     line_spacing=1.1)
pnum(s, 5)

# Slide 6 — Claude
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (5/16)", "Claude (Anthropic)")
bullet_card(s, Inches(0.55), Inches(1.35), Inches(6.0), Inches(2.5), "공식 강점",
    ["Opus 4.6/4.7, Sonnet 4.6: 복잡 추론·장기 에이전트 코딩",
     "맥락 최대 100만 토큰 (Opus/Sonnet 4.6)",
     "Claude Code, Agent SDK, Skills, 구조화 출력",
     "코드 실행·웹 검색·PDF 입력"], BLUE)
bullet_card(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(2.5), "사용자 평가 (참고)",
    ["프로덕션 코딩 품질·가독성 우위 평가 다수",
     "멀티파일 리팩터링·레거시 코드 이해",
     "규제·보안 민감 업무 선호 사례"], ACCENT)
text(s, Inches(0.7), Inches(4.1), Inches(12), Inches(0.4), [[("HW 예시", 13, NAVY, True)]])
text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(1.2),
     [[("\u2022 I2C/SPI 드라이버·HAL 코드 초안, MISRA-C 스타일 점검", 12, GRAY, False)],
      [("\u2022 옛날 펌웨어 코드 설명·인수인계 문서 초안", 12, GRAY, False)]],
     line_spacing=1.1)
pnum(s, 6)

# Slide 7 — 기능 비교 매트릭스
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (6/16)", "기능 비교 매트릭스")
table_simple(
    s, Inches(0.45), Inches(1.5),
    ["항목", "ChatGPT", "Gemini", "Claude"],
    [
        ["코딩·추론", "◎", "◎", "◎"],
        ["에이전트·MCP", "◎", "◎", "◎"],
        ["긴 문서·맥락", "○", "◎", "◎"],
        ["검색·최신정보", "○", "◎", "△"],
        ["멀티모달·OCR", "○", "◎", "○"],
        ["Google 연동", "△", "◎", "△"],
        ["코드 품질(참고)", "○", "○", "◎"],
    ],
    col_w=[Inches(2.5), Inches(3.2), Inches(3.2), Inches(3.2)],
    row_h=Inches(0.58), fsize=12, hsize=11.5
)
text(s, Inches(0.55), Inches(6.35), Inches(12), Inches(0.5),
     [[("◎ 강함  ○ 보통  △ 상대적 약함  |  '코드 품질'은 웹 사용기 참고, 공식 벤치 아님", 10.5, LGRAY, False)]],
     align=PP_ALIGN.CENTER)
pnum(s, 7)

# Slide 8 — 유저 평가 조사 방법
s = add_slide(); set_bg(s, WHITE)
header(s, "유저 평가 기반 비교  (7/16)", "Reddit · X · GitHub · Slack · Discord — 조사 방법")
section(s, Inches(0.55), Inches(1.2), Inches(6), "수집 범위 (5개 채널)")
text(s, Inches(0.7), Inches(1.6), Inches(5.8), Inches(2.4),
     [[("Reddit: 비교·불만 스레드, 고득표(수천 upvotes) 글", 12, GRAY, False)],
      [("X: 개발자·빌더 커뮤니티, 2025~26 Claude 전환 논의", 12, GRAY, False)],
      [("GitHub: Gist·claude-vs-codex 리포, 공식 Issues", 12, GRAY, False)],
      [("Slack: 팀 협업·요약·엔터프라이즈 연동 사용기", 12, GRAY, False)],
      [("Discord: 코딩 커뮤니티 서버·원격 제어·설문 배포", 12, GRAY, False)]],
     line_spacing=1.28)
box(s, Inches(6.75), Inches(1.35), Inches(6.05), Inches(2.6), fill=PEACH, line=ACCENT, line_w=1.2, round_=True)
text(s, Inches(6.95), Inches(1.5), Inches(5.65), Inches(2.3),
     [[("주의 (한계)", 14, ACCENT, True)],
      [("· API로 전체 수집한 것이 아님", 12, GRAY, False)],
      [("· 개인 경험·버전·요금제에 따라 다름", 12, GRAY, False)],
      [("· 홍보·제휴 글 혼재 가능", 12, GRAY, False)],
      [("", 6, GRAY, False)],
      [("→ '반복되는 패턴'만 정리", 12, RED, True)]],
     line_spacing=1.05, space_after=2)
section(s, Inches(0.55), Inches(4.2), Inches(12), "유저들이 공통으로 말하는 것")
text(s, Inches(0.7), Inches(4.65), Inches(12), Inches(2.2),
     [[("\u2022 '최고 하나'보다 업무별로 여러 AI를 스택으로 쓴다", 12.5, NAVY, True)],
      [("\u2022 공식 벤치보다 '매일 쓸 때 한도·속도·실수'가 더 중요하다", 12, GRAY, False)],
      [("\u2022 코딩은 Claude, 긴 문서·Google 업무는 Gemini, 만능·음성은 ChatGPT", 12, GRAY, False)]],
     line_spacing=1.12)
pnum(s, 8)

# Slide 9 — 플랫폼별 유저 목소리
s = add_slide(); set_bg(s, WHITE)
header(s, "유저 평가 기반 비교  (8/16)", "플랫폼별에서 들리는 목소리")
table_simple(
    s, Inches(0.55), Inches(1.4),
    ["플랫폼", "ChatGPT", "Gemini", "Claude"],
    [
        ["Reddit", "만능 칼: 음성·이미지·검색", "Google 앱 통합 편함", "글·코드 품질, 한도 불만"],
        ["X·개발자", "Codex 속도·위임", "Workspace·긴 문서", "빌더층 이동·Claude Code"],
        ["GitHub", "Codex 샌드박스·비용", "무료 티어 보완용", "hooks·plugins·코드 품질"],
        ["Slack", "RAG 연결 비서", "Workspace·Meet 요약", "긴 문서·정책 분석"],
        ["Discord", "빠른 질문·이미지", "설문·토론 활발", "MCP 원격 코딩 제어"],
    ],
    col_w=[Inches(1.5), Inches(3.5), Inches(3.5), Inches(3.75)],
    row_h=Inches(0.6), fsize=11, hsize=12
)
box(s, Inches(0.55), Inches(4.85), Inches(12.25), Inches(1.85), fill=SKY, line=LINE, line_w=1, round_=True)
text(s, Inches(0.8), Inches(4.98), Inches(11.8), Inches(1.6),
     [[("대표 인용 패턴 (참고)", 13, BLUE, True)],
      [("ChatGPT: \"Swiss army knife\" — 빠르고 넓지만 프로덕션 코드는 Claude가 낫다는 의견", 11.5, GRAY, False)],
      [("Gemini: \"검색 레이어\" vs \"긴 문서는 최고\" — 코딩은 평가 갈림", 11.5, GRAY, False)],
      [("Claude: \"scalpel\" — 멀티파일·리팩터링 호평, 느리고 한도 걸린다는 불만", 11.5, GRAY, False)]],
     line_spacing=1.06, space_after=2)
pnum(s, 9)

# Slide 9b — 커뮤니티 사용기 5종 (Reddit·X·GitHub·Slack·Discord)
s = add_slide(); set_bg(s, WHITE)
header(s, "유저 평가 기반 비교  (9/16)", "커뮤니티 사용기 (Reddit·X·GitHub·Slack·Discord)")
table_simple(
    s, Inches(0.5), Inches(1.35),
    ["커뮤니티", "대표 사용기 · 인용 패턴 (참고)"],
    [
        ["Reddit", "\"만능 칼(ChatGPT) vs 정밀 메스(Claude)\"; GPT-5 '답 짧다' 불만 4,600+ upvotes; Claude 한도 불만 최다"],
        ["X·개발자", "2025~26 빌더·개발자층 Claude 이동; Codex(속도·비용·위임) vs Claude Code(코드 품질) 논쟁"],
        ["GitHub", "Gist 비교: Claude 확장성·맥락 / Codex 샌드박스 보안; claude-code·gemini-cli Issues에 한도·자율과잉 불만"],
        ["Slack", "Slack AI: 스레드 요약·액션아이템; ChatGPT를 RAG '프론트도어 비서'로; 'Slack \u2192 작업 자동화' 비교 기준"],
        ["Discord", "Claude Code를 MCP로 원격 제어(폰으로 코딩 지시·수신); 도구 비교·개발자 설문 토론 활발"],
    ],
    col_w=[Inches(1.6), Inches(10.65)], row_h=Inches(0.68), fsize=10.5, hsize=11.5
)
box(s, Inches(0.5), Inches(5.45), Inches(12.3), Inches(1.25), fill=NAVY, round_=True)
text(s, Inches(0.75), Inches(5.58), Inches(11.85), Inches(1.0),
     [[("HW팀 시사점", 13.5, RGBColor(0x9F,0xC4,0xE8), True)],
      [("\u2022 공통 결론: '하나만'이 아니라 업무별 스택 사용 — Reddit·X·GitHub 모두 도구 갈아타기를 권장", 11.5, WHITE, False)],
      [("\u2022 팀 활용은 Slack 요약부터, Discord·MCP 원격 제어는 사내 보안 검토 후 (외부 연동 주의)", 11.5, WHITE, False)]],
     line_spacing=1.08, space_after=2)
pnum(s, 10)

# Slide 9c — 정량 데이터 (설문·통계) (신규)
s = add_slide(); set_bg(s, WHITE)
header(s, "유저 평가 기반 비교  (10/16)", "정량 데이터 (제3자 설문·통계, 참고)")
table_simple(
    s, Inches(0.5), Inches(1.45),
    ["지표", "수치", "출처"],
    [
        ["'가장 사랑하는' 코딩툴", "Claude Code 46% · Cursor 19% · Copilot 9%", "JetBrains'26"],
        ["1차 사용 도구 점유율", "Claude Code 28% · Cursor 24% · Copilot 17% · Codex 11%", "digitalapplied"],
        ["전체 사용(any-use)", "Copilot 58% · Claude Code 54% · Cursor 49%", "digitalapplied"],
        ["멀티툴 동시 사용", "개발자 70~73%가 2개 이상 사용", "Awesome/Ivern"],
        ["멀티에이전트 시간 절감", "11.4h/주 (단일툴 5.2h)", "Ivern 312명"],
        ["AI 도입 vs 신뢰", "도입 84% · 신뢰 29% (전년比 -11p)", "Stack Overflow'25"],
    ],
    col_w=[Inches(3.0), Inches(7.0), Inches(2.3)], row_h=Inches(0.58), fsize=11, hsize=11.5
)
box(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.2), fill=PEACH, line=ACCENT, line_w=1.2, round_=True)
text(s, Inches(0.75), Inches(5.62), Inches(11.8), Inches(1.0),
     [[("읽는 법 (주의)", 13, ACCENT, True)],
      [("표본·방법이 서로 다른 제3자 설문 \u2192 '순위'보다 '경향'으로 해석. ", 11.5, GRAY, False)],
      [("'GPT-5 horrible' 4,600 upvotes, 해지 1.5M 등 언론 인용 수치는 공식 미검증.", 11.5, RED, False)]],
     line_spacing=1.08, space_after=2)
pnum(s, 11)

# Slide 9d — AI별 실사용 불만 상세 (신규)
s = add_slide(); set_bg(s, WHITE)
header(s, "유저 평가 기반 비교  (11/16)", "AI별 실사용 불만 (상세)")
bullet_card(s, Inches(0.45), Inches(1.4), Inches(4.05), Inches(3.6), "ChatGPT",
    ["GPT-5 전환 후 '답 짧음·개성↓·다단계 추론↓' 대규모 불만",
     "Plus 사용 한도(주당 thinking 메시지) 체감",
     "피크시간 품질 변동",
     "안전필터 과도한 거절"], BLUE)
bullet_card(s, Inches(4.65), Inches(1.4), Inches(4.05), Inches(3.6), "Gemini",
    ["'코드 바꾸지 말고 분석만' 지시 무시 등 지시 불이행",
     "자율 과잉: 340파일 변경·28,745줄 삭제 + 허위 점검 로그 사건",
     "멀티스텝에서 워크플로 상태 상실",
     "소비자용 UX(파일 생성·검색 제어 제약)"], ACCENT)
bullet_card(s, Inches(8.85), Inches(1.4), Inches(4.0), Inches(3.6), "Claude",
    ["Rate limit 최대 불만: 주간 캡 1~2일 소진·경고 없이 차단",
     "8개 사용량 버킷 중 하나만 100%여도 throttle",
     "claude.ai·Code·Design 한도 통합 \u2192 한 곳 쓰면 다른 곳↓",
     "상대적으로 느리고 비쌈"], RED)
box(s, Inches(0.45), Inches(5.25), Inches(12.4), Inches(1.45), fill=NAVY, round_=True)
text(s, Inches(0.7), Inches(5.4), Inches(11.9), Inches(1.15),
     [[("HW팀 시사점", 14, RGBColor(0x9F,0xC4,0xE8), True)],
      [("사내 Gauss의 '분당 3~4회 제한'처럼, 외부 AI도 '사용 한도'가 실무 병목이 됩니다.", 12, WHITE, False)],
      [("\u2192 PoC에서 한도·속도·코드 신뢰도(틀린 코드 비율)를 반드시 직접 측정하세요.", 12, WHITE, False)]],
     line_spacing=1.1, space_after=2)
pnum(s, 12)

# Slide 10 — 유저 평가 업무별 비교
s = add_slide(); set_bg(s, WHITE)
header(s, "유저 평가 기반 비교  (12/16)", "업무별 유저 평가 (HW 관점)")
table_simple(
    s, Inches(0.55), Inches(1.48),
    ["업무", "ChatGPT", "Gemini", "Claude"],
    [
        ["코딩 품질", "빠른 초안·Codex", "repo 통째·속도", "가독성·멀티파일 ◎"],
        ["디버깅", "빠른 수정", "보안·맥락 설명", "원인 설명 ◎"],
        ["긴 데이터시트", "보통", "긴 PDF ◎", "긴 글·구조"],
        ["에이전트", "Codex PR·AgentKit", "Antigravity", "Claude Code ◎"],
        ["일상·멀티", "음성·이미지 ◎", "Android·Google ◎", "앱에서는 약함"],
    ],
    col_w=[Inches(2.2), Inches(3.2), Inches(3.2), Inches(3.55)],
    row_h=Inches(0.62), fsize=11, hsize=11.5
)
text(s, Inches(0.55), Inches(5.05), Inches(12.25), Inches(0.45),
     [[("◎ = 유저 후기에서 상대적 강점으로 자주 언급 (공식 벤치 아님)", 10.5, LGRAY, False)]],
     align=PP_ALIGN.CENTER)
box(s, Inches(0.55), Inches(5.55), Inches(12.25), Inches(1.35), fill=NAVY, round_=True)
text(s, Inches(0.8), Inches(5.7), Inches(11.8), Inches(1.05),
     [[("유저 추천 스택 (참고)", 14, RGBColor(0x9F,0xC4,0xE8), True)],
      [("Claude Code(코딩) + Gemini(긴 사양서·OCR) + ChatGPT(만능·음성·빠른 질문)", 12, WHITE, False)],
      [("HW팀 PoC: 업무별로 2주씩 써 보고 한도·속도·실수율을 직접 측정 권장", 12, WHITE, False)]],
     line_spacing=1.1, space_after=2)
pnum(s, 13)

# Slide 11 — 공식 vs 유저
s = add_slide(); set_bg(s, WHITE)
header(s, "유저 평가 기반 비교  (13/16)", "공식 소개 vs 유저 평가 — 차이")
table_simple(
    s, Inches(0.55), Inches(1.5),
    ["AI", "공식이 말하는 것", "유저가 추가로 말하는 것"],
    [
        ["ChatGPT", "GPT-5.x, AgentKit, 생태계", "만능·Codex 빠름, 한도·품질 변동"],
        ["Gemini", "3.5 Flash 벤치, Google 연동", "Workspace 편함, 코딩은 기대↓ 의견도"],
        ["Claude", "Opus/Sonnet, 1M, Claude Code", "코드 1위 인상, 한도·속도 병목"],
    ],
    col_w=[Inches(1.35), Inches(5.2), Inches(5.7)], row_h=Inches(0.95), fsize=11.5, hsize=12
)
box(s, Inches(0.55), Inches(5.15), Inches(12.25), Inches(1.75), fill=PEACH, line=ACCENT, line_w=1.2, round_=True)
text(s, Inches(0.8), Inches(5.3), Inches(11.8), Inches(1.45),
     [[("도입 시 체크리스트", 14, ACCENT, True)],
      [("1. 공식 문서로 '기능이 되는지' 확인", 12, GRAY, False)],
      [("2. 유저 후기로 '매일 쓸 때 불편한지' PoC에서 확인 (한도·속도·코드 검증)", 12, GRAY, False)],
      [("3. 사내 Gauss와 역할 분리 — 외부는 보안 승인 후 고급 보조", 12, GRAY, False)]],
     line_spacing=1.1, space_after=2)
pnum(s, 14)

# Slide 12 — HW팀 시사점 + 선택 가이드
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (14/16)", "HW 개발팀 적용 및 선택 가이드")
table_simple(
    s, Inches(0.55), Inches(1.5),
    ["업무", "우선 검토", "이유 (쉬운 말)"],
    [
        ["펌웨어·드라이버 초안", "Claude, ChatGPT", "코드 구조·주석 품질 평가 높음"],
        ["긴 데이터시트 요약", "Gemini", "긴 문서·검색이 공식 강점"],
        ["측정 로그·표 정리", "3종 모두", "반복 정리는 모두 활용 가능"],
        ["스캔 PDF·회로도 OCR", "Gemini", "문서·이미지 추론 공식 사례"],
        ["다단계 자동화", "ChatGPT, Claude", "AgentKit / Claude Code"],
        ["최종 회로·양산 판단", "사람", "AI는 보조, 실측 필수"],
    ],
    col_w=[Inches(2.8), Inches(2.5), Inches(6.95)], row_h=Inches(0.55), fsize=11, hsize=11.5
)
box(s, Inches(0.55), Inches(5.35), Inches(12.25), Inches(1.35), fill=SKY, line=LINE, line_w=1, round_=True)
text(s, Inches(0.8), Inches(5.5), Inches(11.8), Inches(1.05),
     [[("선택 가이드", 14, BLUE, True)],
      [("코딩 중심 \u2192 Claude  |  긴 사양서·검색 \u2192 Gemini  |  범용·도구 연결 \u2192 ChatGPT", 12, GRAY, False)],
      [("복수 도입·역할 분담 권장 (예: Claude 코딩 + Gemini 문서)", 12, GRAY, False)]],
     line_spacing=1.1, space_after=2)
pnum(s, 15)

# Slide 13 — 결론
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (15/16)", "결론 및 제언")
cols = [
    ("당장 PoC 가능", GREEN, ["코드·문서 초안", "측정 데이터 정리", "로그 분석", "반복 업무 자동화"]),
    ("신중 적용", ACCENT, ["회로·부품 최종 판단", "불량 원인 확정", "양산 영향 평가"]),
    ("개선·검토 과제", RED, ["보안·법무 승인", "사내 Gauss와 역할 분리", "팀 공통 가이드", "도구별 라이선스"]),
]
left = Inches(0.65); top = Inches(1.45); cw = Inches(3.95); gap = Inches(0.25); ch = Inches(3.2)
for i, (title, color, items) in enumerate(cols):
    x = left + i * (cw + gap)
    box(s, x, top, cw, ch, fill=WHITE, line=color, line_w=1.8, round_=True)
    box(s, x, top, cw, Inches(0.55), fill=color, round_=True)
    text(s, x, top, cw, Inches(0.55), [[(title, 14, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    yy = top + Inches(0.75)
    for item in items:
        text(s, x + Inches(0.25), yy, cw - Inches(0.5), Inches(0.35), [[("- " + item, 11.5, GRAY, False)]])
        yy += Inches(0.5)
box(s, Inches(0.65), Inches(4.95), Inches(12.05), Inches(1.55), fill=NAVY, round_=True)
text(s, Inches(0.9), Inches(5.1), Inches(11.5), Inches(1.25),
     [[("최종 결론", 15, RGBColor(0x9F,0xC4,0xE8), True)],
      [("외부 AI는 사내 Gauss를 대체하기보다, 보안 검토 후 '고급 보조'로 쓰는 것이 현실적입니다.", 12.5, WHITE, False)],
      [("6월 도입 전: 업무별 PoC 1~2건 선정 \u2192 효과·보안 측정 \u2192 팀 표준 가이드화를 제안합니다.", 12.5, WHITE, False)]],
     line_spacing=1.1, space_after=2)
pnum(s, 16)

# Slide 14 — 출처
s = add_slide(); set_bg(s, WHITE)
header(s, "외부 AI 3종 비교  (16/16)", "출처")
sources = [
    ("공식", "OpenAI / Google Gemini / Anthropic Claude 문서"),
    ("Reddit", "Foxafox 요약, RoboRhythms, TechRadar, ClaudeMeter, OpenTools"),
    ("GitHub", "Haseeb Qureshi Gist, rohittcodes/claude-vs-codex, claude-code·gemini-cli Issues, arXiv 20,574 세션"),
    ("설문", "JetBrains'26, digitalapplied, Ivern(312명), Stack Overflow'25"),
    ("Slack·Discord", "editorialge, blockchain-council, State of AI Coding 2026"),
    ("X·참고", "Vapvarun, CatDoes, Medium 2026 field report"),
    ("문서", "외부AI_3종_비교조사.md + 유저평가_비교조사.md"),
]
y = Inches(1.55)
for tag, desc in sources:
    box(s, Inches(0.7), y, Inches(1.1), Inches(0.42), fill=BLUE if tag == "공식" else ACCENT, round_=True)
    text(s, Inches(0.7), y, Inches(1.1), Inches(0.42), [[(tag, 11, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, Inches(2.0), y, Inches(10.5), Inches(0.42), [[(desc, 12, GRAY, False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)
text(s, Inches(0.7), Inches(5.5), Inches(12), Inches(1.2),
     [[("상세 URL은 동봉 마크다운 '외부AI_3종_비교조사.md' 참조", 11, LGRAY, False)],
      [("모델명·기능은 도입 시점에 공식 페이지에서 재확인 필요", 11, LGRAY, False)]],
     line_spacing=1.1)
pnum(s, 17)

prs.save("외부AI_3종_비교조사.pptx")
print("saved:", len(prs.slides._sldIdLst), "slides")
