# -*- coding: utf-8 -*-
"""NC HW 개발팀 — Cline SR 사내 활용 (표지·도구 설명 + 활용 범위·한계·대조)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x14, 0x2A, 0x4A)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
SKY    = RGBColor(0xE8, 0xF1, 0xFA)
GRAY   = RGBColor(0x44, 0x4A, 0x55)
LGRAY  = RGBColor(0x8A, 0x92, 0x9E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF2, 0x7E, 0x2E)
GREEN  = RGBColor(0x27, 0xAE, 0x60)
RED    = RGBColor(0xC0, 0x39, 0x2B)
PEACH  = RGBColor(0xFE, 0xF1, 0xE6)
LINE   = RGBColor(0xD5, 0xDD, 0xE6)
FONT   = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide(): return prs.slides.add_slide(BLANK)
def set_bg(slide, color):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, l, t, w, h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb

def header(slide, kicker, title):
    box(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    box(slide, 0, Inches(1.0), SW, Pt(3), fill=ACCENT)
    text(slide, Inches(0.55), Inches(0.13), Inches(12), Inches(0.3),
         [[(kicker, 11, RGBColor(0x9F,0xC4,0xE8), True)]])
    text(slide, Inches(0.55), Inches(0.4), Inches(12.2), Inches(0.55),
         [[(title, 23, WHITE, True)]])

def section(slide, l, t, w, label, color=BLUE):
    box(slide, l, t, Pt(5), Inches(0.32), fill=color)
    text(slide, l + Inches(0.12), t - Inches(0.02), w, Inches(0.36),
         [[(label, 14.5, NAVY, True)]])

def pnum(slide, n):
    text(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
         [[(str(n), 10, LGRAY, False)]], align=PP_ALIGN.RIGHT)

def table_simple(slide, l, t, headers, rows, col_w, row_h, head_fill=BLUE,
                 fsize=10.5, hsize=11):
    x = l
    for j, htxt in enumerate(headers):
        cw = col_w[j]
        box(slide, x, t, cw, row_h, fill=head_fill, line=WHITE, line_w=1)
        text(slide, x, t, cw, row_h, [[(htxt, hsize, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cw
    y = t + row_h
    for ri, row in enumerate(rows):
        x = l; fill = WHITE if ri % 2 == 0 else SKY
        for j, cell in enumerate(row):
            cw = col_w[j]
            box(slide, x, y, cw, row_h, fill=fill, line=LINE, line_w=0.75)
            text(slide, x + Pt(3), y, cw - Pt(6), row_h, [[(cell, fsize, GRAY, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.92)
            x += cw
        y += row_h
    return y

KICKER = "NC HW 개발팀  |  Cline SR 사내 활용"
TOTAL = 6

def tool_card(slide, l, t, w, h, name, desc, scope_lines, color=BLUE):
    box(slide, l, t, w, h, fill=WHITE, line=color, line_w=1.5, round_=True)
    box(slide, l, t, w, Inches(0.42), fill=color, round_=True)
    text(slide, l, t, w, Inches(0.42), [[(name, 13, WHITE, True)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(slide, l + Inches(0.15), t + Inches(0.5), w - Inches(0.3), Inches(0.55),
         [[(desc, 10.2, GRAY, False)]], line_spacing=1.05)
    text(slide, l + Inches(0.15), t + Inches(1.05), w - Inches(0.3), Inches(0.22),
         [[("활용 범위", 10.5, color, True)]])
    yy = t + Inches(1.28)
    for line in scope_lines:
        text(slide, l + Inches(0.18), yy, w - Inches(0.32), Inches(0.28),
             [[("\u2022 " + line, 9.8, GRAY, False)]], line_spacing=1.0)
        yy += Inches(0.3)

# ============================================================
# Slide 0 — 표지
# ============================================================
s = add_slide(); set_bg(s, NAVY)
box(s, 0, Inches(3.9), SW, Pt(4), fill=ACCENT)
text(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
     [[("Cline SR 사내 활용", 40, WHITE, True)]],
     line_spacing=1.1)
text(s, Inches(0.92), Inches(4.35), Inches(11.5), Inches(0.9),
     [[("NC HW 개발팀  |  Gauss · Ollama · GitHub Enterprise 연계", 16, RGBColor(0xC9,0xD8,0xEA), False)],
      [("VS Code 기반 사내 AI 코딩 환경 활용 방안", 14, LGRAY, False)]],
     line_spacing=1.15)
pnum(s, 1)

# ============================================================
# Slide 1 — 도구 설명 및 활용 범위
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER + "  (2/6)", "사내 도구 구성 — 설명 및 활용 범위")

cw = Inches(6.05); ch = Inches(2.72); gap_x = Inches(0.2); gap_y = Inches(0.18)
x0 = Inches(0.5); y0 = Inches(1.35)

tool_card(s, x0, y0, cw, ch, "Cline SR",
    "VS Code 확장 AI 코딩 에이전트. 지시 → 계획 → 실행으로 코드 읽기·작성·수정, 터미널 명령 수행.",
    ["펌웨어·드라이버·HAL 코드 초안·리팩터링", "멀티파일 일괄 수정·디버깅 보조",
     "측정·검증 Python 스크립트 생성", "레거시 코드 분석·설명"], BLUE)

tool_card(s, x0 + cw + gap_x, y0, cw, ch, "Gauss",
    "사내 대형 언어모델(LLM) API. Cline SR의 두뇌 역할. 사내망 내 처리로 외부 유출 없음.",
    ["일반 코딩·문서·코드 리뷰 질의", "시험·이슈·회의 보고서 초안",
     "데이터시트 기반 코드·설명 생성", "분당 3~4회 호출 제한 → 장시간 작업 시 대기"], ACCENT)

tool_card(s, x0, y0 + ch + gap_y, cw, ch, "Ollama",
    "개발자 PC에서 구동하는 로컬 LLM. 완전 오프라인·망분리 환경에서 Cline SR과 연동.",
    ["고민감·망분리 구역 자료 처리", "간단한 코드·문서 보조 (로컬)",
     "GPU 없는 PC는 응답 속도 현저히 저하", "대용량·장문 분석은 Gauss 권장"], GREEN)

tool_card(s, x0 + cw + gap_x, y0 + ch + gap_y, cw, ch, "GitHub Enterprise",
    "사내 Git 저장소·협업 플랫폼. 코드 이력·브랜치·PR·이슈·코드 검색·리뷰 워크플로.",
    ["소스코드 버전 관리·브랜치 전략", "PR·코드 리뷰·이슈 트래킹",
     "팀 내 코드 검색·변경 이력 추적", "Cline SR 작업 결과물의 사내 저장·공유"], NAVY)

pnum(s, 2)

# ============================================================
# Slide 2 — HW 개발 업무 전체 활용 범위
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER + "  (3/6)", "HW 개발 업무 — 사내 AI 환경 활용 범위")

text(s, Inches(0.55), Inches(1.15), Inches(12.2), Inches(0.38),
     [[("목표: 정형·반복 업무를 AI가 직접 수행하도록 확대하고, 엔지니어는 설계·검증·판단에 집중", 11.5, GRAY, False)]])

table_simple(
    s, Inches(0.45), Inches(1.58),
    ["업무 영역", "활용 내용", "주요 도구"],
    [
        ["펌웨어·드라이버", "C/HAL·I2C/SPI/UART 코드 초안·리팩터링", "VS Code + Cline SR + Gauss"],
        ["레지스터·설정", "데이터시트 기반 헤더·초기화 코드 생성", "Cline SR + Gauss"],
        ["측정·검증 자동화", "계측기 제어·로그 수집 Python 스크립트", "Cline SR + Gauss / Ollama"],
        ["데이터·로그 분석", "측정값 파싱·표·그래프·이상 패턴 정리", "Cline SR + Gauss"],
        ["문서·보고", "시험·이슈·회의·인수인계 문서 초안", "Gauss / Cline SR"],
        ["코드 리뷰·품질", "MISRA-C 스타일·예외처리·위험 구간 점검", "Cline SR + Gauss"],
        ["저장소·협업", "브랜치·PR·이슈·코드 검색·리뷰 워크플로", "GitHub Enterprise"],
        ["레거시·HDL 보조", "구형 펌웨어 설명·testbench 초안", "Cline SR + Gauss / Ollama"],
    ],
    col_w=[Inches(2.35), Inches(6.35), Inches(4.15)],
    row_h=Inches(0.52), fsize=10.2, hsize=11
)

box(s, Inches(0.45), Inches(5.95), Inches(12.4), Inches(1.05), fill=NAVY, round_=True)
text(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(1.05),
     [[("환경 구성 요약", 13, RGBColor(0x9F,0xC4,0xE8), True)],
      [("VS Code(IDE) + Cline SR(에이전트) + Gauss(사내 LLM) / Ollama(로컬 LLM) + GitHub Enterprise(사내 저장소)", 11.5, WHITE, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.12, space_after=2)
pnum(s, 3)

# ============================================================
# Slide 3 — 지금 당장 활용 가능한 분야
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER + "  (4/6)", "지금 당장 활용 가능한 분야")

section(s, Inches(0.55), Inches(1.18), Inches(12), "즉시 적용 가능 (보안 검토 완료 환경 기준)")
table_simple(
    s, Inches(0.55), Inches(1.58),
    ["분야", "활용 방법", "HW 개발 예시"],
    [
        ["코드 초안", "Cline SR이 파일 단위로 작성·수정", "센서·PMIC·충전 IC 드라이버·HAL 초안"],
        ["레지스터 맵", "표 기반 헤더·초기화 코드 생성", "데이터시트 레지스터 정의 → .h 변환"],
        ["측정 스크립트", "Python·SCPI 자동화 코드", "전압·전류·온도·소비전력 측정 루프"],
        ["로그·데이터 정리", "파싱·표·요약 문장 생성", "시험 로그·CSV → 보고용 표·그래프"],
        ["문서 초안", "코드·이슈 기반 문서화", "시험 결과·회의록·개발 가이드 초안"],
        ["코드 점검", "규칙·패턴·누락 검토", "예외처리·매직넘버·중복 코드 지적"],
        ["Git 워크플로", "PR·이슈·코드 검색", "브랜치 관리·리뷰·이력 추적 (GitHub Enterprise)"],
    ],
    col_w=[Inches(1.55), Inches(3.15), Inches(7.55)],
    row_h=Inches(0.58), fsize=11, hsize=11.5
)

yb = Inches(5.35)
box(s, Inches(0.55), yb, Inches(6.0), Inches(1.35), fill=SKY, line=LINE, line_w=1, round_=True)
text(s, Inches(0.78), yb, Inches(5.55), Inches(1.35),
     [[("권장 운영", 13, BLUE, True)],
      [("\u2022 Gauss: 일반 코딩·문서 (고성능)", 11.5, GRAY, False)],
      [("\u2022 Ollama: 망분리·고민감 자료 (로컬)", 11.5, GRAY, False)],
      [("\u2022 긴 작업은 짧은 단위로 분할 요청", 11.5, GRAY, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=2)
box(s, Inches(6.75), yb, Inches(6.05), Inches(1.35), fill=NAVY, round_=True)
text(s, Inches(6.98), yb, Inches(5.6), Inches(1.35),
     [[("기대 효과", 13, RGBColor(0x9F,0xC4,0xE8), True)],
      [("반복·정형 업무 시간 단축", 11.5, WHITE, False)],
      [("개발자는 설계·실측·검증에 집중", 11.5, WHITE, False)],
      [("사내 망 내 데이터 유출 없음", 11.5, WHITE, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1, space_after=2)
pnum(s, 4)

# ============================================================
# Slide 4 — 현재 사내 조건의 한계
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER + "  (5/6)", "현재 사내 조건의 한계")

L = Inches(0.5); W = Inches(12.35)
limits = [
    ("외부 MCP·API 사용 제한",
     "외부 문서 검색·외부 저장소·외부 분석 도구 자동 연동 불가. Cline SR이 스스로 최신 공개 자료·외부 API를 호출할 수 없음."),
    ("임베딩·청킹 성능 제한",
     "장문 데이터시트·설계서의 의미 기반 검색·분할 처리 품질이 낮음. 수백 페이지 사양서에서 특정 설정값을 정확히 찾기 어려움."),
    ("OCR 성능 제한",
     "스캔 PDF·회로도 캡처·이미지 내 텍스트 인식 품질이 낮음. 도면·스캔 자료를 바로 구조화·분석하기 어려움."),
    ("Gauss API 호출 제한",
     "분당 약 3~4회 호출 제한. Cline SR의 다단계·다회 수정 작업 시 대기 시간이 누적되어 작업 속도 저하."),
    ("Ollama 로컬 성능 제한",
     "개인 PC에 GPU가 없으면 응답·추론 속도가 현저히 느림. 대용량 코드·장문 분석은 실무 적용이 어려움."),
]

y = Inches(1.52)
for idx, (title, desc) in enumerate(limits):
    bh = Inches(0.88)
    box(s, L, y, W, bh, fill=PEACH, line=ACCENT, line_w=1.0, round_=True)
    box(s, L, y, Pt(6), bh, fill=ACCENT)
    text(s, L + Inches(0.22), y + Inches(0.1), Inches(3.15), Inches(0.32), [[(title, 12, NAVY, True)]])
    text(s, L + Inches(3.45), y + Inches(0.1), Inches(8.75), Inches(0.68),
         [[(desc, 10.8, GRAY, False)]], line_spacing=1.05)
    y += bh + Inches(0.1)

box(s, L, Inches(6.72), W, Inches(0.48), fill=NAVY, round_=True)
text(s, L, Inches(6.72), W, Inches(0.48),
     [[("핵심: 보안·망분리 환경에서는 유용하나, 외부 연동·고급 문서/OCR·고속 대량 호출은 현재 구조상 제약", 11.2, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
pnum(s, 5)

# ============================================================
# Slide 5 — 사내 vs 외부 도구 대조 및 결론
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, KICKER + "  (6/6)", "사내 AI 환경 vs 외부 도구 — 대조 및 결론")

# 좌: 사내 / 우: 외부
half_w = Inches(6.05)
gap = Inches(0.2)
lx = Inches(0.5); rx = lx + half_w + gap; ty = Inches(1.35); th = Inches(4.85)

box(s, lx, ty, half_w, th, fill=SKY, line=BLUE, line_w=2, round_=True)
box(s, lx, ty, half_w, Inches(0.5), fill=BLUE, round_=True)
text(s, lx, ty, half_w, Inches(0.5), [[("현재 사내 환경 (가능)", 14, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
internal_items = [
    "펌웨어·드라이버·HAL 코드 초안·수정",
    "측정·검증 Python 스크립트 자동화",
    "로그·측정 데이터 파싱·표·요약",
    "시험·이슈·회의 문서 초안",
    "사내 코드·GitHub Enterprise 기반 리뷰",
    "망분리·로컬(Ollama) 고민감 자료 처리",
    "Gauss·Cline SR 기반 다단계 코드 작업 (호출 한도 내)",
]
yy = ty + Inches(0.62)
for item in internal_items:
    text(s, lx + Inches(0.25), yy, half_w - Inches(0.45), Inches(0.38),
         [[("\u2713 " + item, 11.2, GRAY, False)]], line_spacing=1.0)
    yy += Inches(0.52)

box(s, rx, ty, half_w, th, fill=PEACH, line=RED, line_w=2, round_=True)
box(s, rx, ty, half_w, Inches(0.5), fill=RED, round_=True)
text(s, rx, ty, half_w, Inches(0.5), [[("외부 도구 (인증·정책 제약)", 14, WHITE, True)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
external_items = [
    "외부 MCP·API·플러그인 자동 연동",
    "ChatGPT·Gemini·Claude 등 외부 LLM 직접 연동",
    "웹·최신 공개 자료 실시간 검색·인용",
    "고품질 임베딩·장문 사양서 의미 검색",
    "고정밀 OCR (스캔 PDF·회로도)",
    "대량·고속 API 호출 (분당 수십~수백 회)",
    "멀티모달·에이전트 클라우드 위임 (Codex 등)",
]
yy = ty + Inches(0.62)
for item in external_items:
    text(s, rx + Inches(0.25), yy, half_w - Inches(0.45), Inches(0.38),
         [[("\u2717 " + item, 11.2, GRAY, False)]], line_spacing=1.0)
    yy += Inches(0.52)

# 중앙 VS
text(s, Inches(6.35), Inches(3.2), Inches(0.65), Inches(0.5),
     [[("VS", 16, ACCENT, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

box(s, Inches(0.5), Inches(6.38), Inches(12.35), Inches(0.88), fill=NAVY, round_=True)
text(s, Inches(0.75), Inches(6.38), Inches(11.85), Inches(0.88),
     [[("결론", 14, RGBColor(0x9F,0xC4,0xE8), True)],
      [("현재 사내 AI 환경에서 활용 가능한 부분은 ", 11.5, WHITE, False),
       ("코드·스크립트 초안, 측정·로그 정리, 문서 초안, 사내 Git 기반 협업, 망분리 로컬 처리", 11.5, WHITE, True),
       (" 이다.", 11.5, WHITE, False)],
      [("다만 사내 조건의 한계로 불가능하거나 어려운 부분은 ", 11.5, WHITE, False),
       ("외부 MCP·API 연동, 고품질 장문 검색·OCR, 대량 고속 호출, 외부 LLM·최신 웹 자료 자동 활용", 11.5, WHITE, True),
       (" 이다.", 11.5, WHITE, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.15, space_after=2)
pnum(s, 6)

prs.save("HW개발팀_AI활용방안_요약.pptx")
print("saved summary:", len(prs.slides._sldIdLst), "slides")
