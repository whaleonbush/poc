# -*- coding: utf-8 -*-
"""HW 개발팀 AI 코딩 에이전트 활용 방안 — 2페이지 요약본"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY   = RGBColor(0x14, 0x2A, 0x4A)
BLUE   = RGBColor(0x1F, 0x6F, 0xB2)
SKY    = RGBColor(0xE8, 0xF1, 0xFA)
GRAY   = RGBColor(0x44, 0x4A, 0x55)
LGRAY  = RGBColor(0x8A, 0x92, 0x9E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0xF2, 0x7E, 0x2E)
LINE   = RGBColor(0xD5, 0xDD, 0xE6)
FONT   = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)

def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def box(slide, l, t, w, h, fill=None, line=None, line_w=1.0, round_=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shp, l, t, w, h)
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def text(slide, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2); tf.margin_top = tf.margin_bottom = Pt(2)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space_after); p.line_spacing = line_spacing
        for (txt, size, color, bold) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = FONT
    return tb

def header(slide, kicker, title):
    box(slide, 0, 0, SW, Inches(1.0), fill=NAVY)
    box(slide, 0, Inches(1.0), SW, Pt(3), fill=ACCENT)
    text(slide, Inches(0.55), Inches(0.13), Inches(12), Inches(0.3),
         [[(kicker, 11, RGBColor(0x9F,0xC4,0xE8), True)]])
    text(slide, Inches(0.55), Inches(0.4), Inches(12.2), Inches(0.55),
         [[(title, 23, WHITE, True)]])

def section(slide, l, t, w, label, color=BLUE):
    box(slide, l, t, Pt(5), Inches(0.32), fill=color)
    text(slide, l + Inches(0.12), t - Inches(0.02), w, Inches(0.36),
         [[(label, 14.5, NAVY, True)]])

def table_simple(slide, l, t, headers, rows, col_w, row_h, head_fill=BLUE,
                 fsize=11, hsize=11.5):
    x = l
    for j, htxt in enumerate(headers):
        cw = col_w[j]
        box(slide, x, t, cw, row_h, fill=head_fill, line=WHITE, line_w=1)
        text(slide, x, t, cw, row_h, [[(htxt, hsize, WHITE, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        x += cw
    y = t + row_h
    for ri, row in enumerate(rows):
        x = l; fill = WHITE if ri % 2 == 0 else SKY
        for j, cell in enumerate(row):
            cw = col_w[j]
            box(slide, x, y, cw, row_h, fill=fill, line=LINE, line_w=0.75)
            al = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            text(slide, x + Pt(4), y, cw - Pt(8), row_h, [[(cell, fsize, GRAY, False)]],
                 align=al, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
            x += cw
        y += row_h
    return y

# ============================================================
# Slide 1 — 개요 & 구성/보안
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "HW 개발팀 AI 코딩 에이전트 활용 방안  (1/2)", "도입 개요 및 구성")

# 좌측: 도입 개요
section(s, Inches(0.55), Inches(1.25), Inches(6), "1. 도입 개요")
ov = [
    ("무엇을", "VS Code 기반 AI 코딩 에이전트 Cline SR을 사내 LLM(Gauss)·로컬 LLM(Ollama)과 연동해 HW 개발 업무에 도입"),
    ("왜", "코드·설계자산을 외부로 보내지 않는 사내망·로컬 환경에서 AI 활용 → 보안 위험 없이 생산성 향상"),
    ("효과", "펌웨어·테스트·문서·리뷰 등 반복·정형 업무 자동화 → 개발 리드타임 단축, 검증 품질 향상"),
]
y = Inches(1.75)
for k, v in ov:
    box(s, Inches(0.55), y, Inches(6.0), Inches(1.25), fill=SKY, line=LINE, line_w=1, round_=True)
    box(s, Inches(0.55), y, Pt(6), Inches(1.25), fill=BLUE)
    text(s, Inches(0.8), y + Inches(0.12), Inches(5.5), Inches(0.35), [[(k, 14, BLUE, True)]])
    text(s, Inches(0.8), y + Inches(0.45), Inches(5.55), Inches(0.75), [[(v, 12.5, GRAY, False)]],
         line_spacing=1.08)
    y += Inches(1.4)

# Gauss vs Ollama 간단 비교 (좌측 하단)
section(s, Inches(0.55), Inches(6.05), Inches(6), "운영 원칙")
text(s, Inches(0.7), Inches(6.5), Inches(6.0), Inches(0.8),
     [[("일반 업무 \u2192 ", 12.5, GRAY, True),("Gauss", 12.5, BLUE, True),("(사내망·고성능)", 11.5, LGRAY, False)],
      [("최고 민감·망분리 \u2192 ", 12.5, GRAY, True),("Ollama", 12.5, ACCENT, True),("(PC 내 완전 오프라인)", 11.5, LGRAY, False)]],
     line_spacing=1.15, space_after=3)

# 우측: 구성/동작 개념도
section(s, Inches(7.0), Inches(1.25), Inches(6), "2. 도구 구성 및 동작")
rx = Inches(7.0)
box(s, rx, Inches(1.75), Inches(5.75), Inches(0.7), fill=NAVY, round_=True)
text(s, rx, Inches(1.75), Inches(5.75), Inches(0.7),
     [[("개발자", 14, WHITE, True),("  · VS Code + Cline SR(에이전트)", 12, RGBColor(0xC9,0xD8,0xEA), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, rx, Inches(2.5), Inches(5.75), Inches(0.3), [[("\u25BC  두뇌(LLM) 선택", 11, LGRAY, True)]],
     align=PP_ALIGN.CENTER)
box(s, rx, Inches(2.85), Inches(2.8), Inches(1.0), fill=BLUE, round_=True)
text(s, rx, Inches(2.85), Inches(2.8), Inches(1.0),
     [[("Gauss 사내 API", 13, WHITE, True)],[("사내망·고성능·팀 공용", 11, RGBColor(0xD8,0xE6,0xF4), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(s, rx + Inches(2.95), Inches(2.85), Inches(2.8), Inches(1.0), fill=BLUE, round_=True)
text(s, rx + Inches(2.95), Inches(2.85), Inches(2.8), Inches(1.0),
     [[("Ollama 로컬 LLM", 13, WHITE, True)],[("PC 내 완전 오프라인", 11, RGBColor(0xD8,0xE6,0xF4), False)]],
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
# 보안 강조
box(s, rx, Inches(4.2), Inches(5.75), Inches(1.05), fill=RGBColor(0xFE,0xF1,0xE6), line=ACCENT, line_w=1.5, round_=True)
text(s, rx + Inches(0.2), Inches(4.2), Inches(5.4), Inches(1.05),
     [[("\U0001F512 핵심 보안 포인트", 13, ACCENT, True)],
      [("입력 코드·결과물이 사내망/로컬 PC 밖으로 전혀", 12, GRAY, False)],
      [("나가지 않음 \u2192 외부 유출 위험 없음", 12, GRAY, False)]],
     anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0, space_after=1)
# 구성요소 한 줄 요약
text(s, rx, Inches(5.45), Inches(5.75), Inches(1.6),
     [[("· VS Code : 표준 IDE(편집기)", 11.5, GRAY, False)],
      [("· Cline SR : 코드 읽기/작성/수정·명령 실행을 자동 수행하는 에이전트", 11.5, GRAY, False)],
      [("· Gauss : 사내 LLM API (두뇌)", 11.5, GRAY, False)],
      [("· Ollama : PC 구동 로컬 LLM (오프라인)", 11.5, GRAY, False)]],
     line_spacing=1.1, space_after=2)
text(s, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3), [[("1", 10, LGRAY, False)]], align=PP_ALIGN.RIGHT)

# ============================================================
# Slide 2 — 활용 방안 & 기대효과 & 제언
# ============================================================
s = add_slide(); set_bg(s, WHITE)
header(s, "HW 개발팀 AI 코딩 에이전트 활용 방안  (2/2)", "활용 방안 및 기대 효과")

section(s, Inches(0.55), Inches(1.2), Inches(8), "3. HW 개발팀 활용 방안")
table_simple(
    s, Inches(0.55), Inches(1.62),
    ["구분", "활용 방안", "내용 / 기대 효과"],
    [
        ["펌웨어", "드라이버·통신 코드 작성", "C/C++ 드라이버·HAL·I2C/SPI/UART/CAN 초안 생성"],
        ["펌웨어", "레지스터 맵 자동화", "데이터시트 → 헤더(.h) 변환, 입력 오류 감소"],
        ["검증", "테스트 자동화 스크립트", "계측기 제어(SCPI)·측정 자동화 Python 생성"],
        ["검증", "측정 데이터 분석", "로그 파싱·그래프/리포트 생성, 분석시간 단축"],
        ["품질", "코드 리뷰·디버깅 보조", "MISRA-C 점검, 버그·예외처리 누락 지적"],
        ["문서", "문서화 자동화", "주석·README·설계/테스트 리포트 초안 작성"],
        ["유지", "레거시 코드 이해", "오래된/인수인계 코드 분석·설명, 온보딩 가속"],
    ],
    col_w=[Inches(1.3), Inches(3.0), Inches(7.45)], row_h=Inches(0.52), fsize=11.5, hsize=12
)

# 우측 하단: 기대효과 + 리스크 + 제언  → 하단 가로 배치
yb = Inches(5.95)
# 기대 효과
box(s, Inches(0.55), yb, Inches(4.0), Inches(1.15), fill=SKY, line=LINE, line_w=1, round_=True)
text(s, Inches(0.75), yb + Inches(0.1), Inches(3.7), Inches(1.0),
     [[("기대 효과", 13, BLUE, True)],
      [("생산성↑·품질↑·지식 자산화", 11.5, GRAY, False)],
      [("보안과 생산성 동시 확보", 11.5, GRAY, False)]], line_spacing=1.1, space_after=2)
# 리스크
box(s, Inches(4.7), yb, Inches(4.0), Inches(1.15), fill=RGBColor(0xFE,0xF1,0xE6), line=ACCENT, line_w=1, round_=True)
text(s, Inches(4.9), yb + Inches(0.1), Inches(3.7), Inches(1.0),
     [[("고려사항", 13, ACCENT, True)],
      [("AI는 보조 도구 → 개발자 검토·", 11.5, GRAY, False)],
      [("실측 검증 후 적용을 원칙으로", 11.5, GRAY, False)]], line_spacing=1.1, space_after=2)
# 제언
box(s, Inches(8.85), yb, Inches(3.95), Inches(1.15), fill=NAVY, round_=True)
text(s, Inches(9.05), yb + Inches(0.1), Inches(3.65), Inches(1.0),
     [[("제언", 13, RGBColor(0x9F,0xC4,0xE8), True)],
      [("저위험·고효용 과제 중심", 11.5, WHITE, False)],
      [("단계적 시범 적용(PoC) 후 확산", 11.5, WHITE, False)]], line_spacing=1.1, space_after=2)
text(s, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3), [[("2", 10, LGRAY, False)]], align=PP_ALIGN.RIGHT)

prs.save("HW개발팀_AI코딩에이전트_활용방안_2page.pptx")
print("saved 2page:", len(prs.slides._sldIdLst))
